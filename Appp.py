import os
import streamlit as st
import pandas as pd
from langchain_groq import ChatGroq
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_core.documents import Document
from dotenv import load_dotenv
from pathlib import Path
import docx
import pytesseract
from PIL import Image, ImageEnhance
import speech_recognition as sr
import re
from datetime import datetime
import pdfkit
import tempfile
import groq
import time
import requests
from bs4 import BeautifulSoup
import uuid
import json
from pptx import Presentation
import logging
import shutil

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configure wkhtmltopdf path for pdfkit
pdfkit_config = pdfkit.configuration(wkhtmltopdf=r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe")

# === [ Load Keys & Setup ] ===
load_dotenv()
groq_api_key = os.getenv("GROQ_API_KEY")
if not groq_api_key:
    st.error("GROQ_API_KEY not found in .env file. Please set a valid key.")
    st.stop()
logger.info(f"GROQ_API_KEY: {groq_api_key[:10]}...")
os.environ['GOOGLE_API_KEY'] = os.getenv("GOOGLE_API_KEY")

embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
try:
    llm = ChatGroq(groq_api_key=groq_api_key, model_name="llama-3.1-8b-instant")
except Exception as e:
    st.error(f"Failed to initialize Groq LLM: {e}. Please check your Groq API key.")
    st.stop()

STORAGE_PATH = "./uploaded_files"
VECTOR_STORE_PATH = "./faiss_index"
TRAINING_DATA_PATH = "./training_data"

os.makedirs(STORAGE_PATH, exist_ok=True)
os.makedirs(VECTOR_STORE_PATH, exist_ok=True)
os.makedirs(TRAINING_DATA_PATH, exist_ok=True)

# Session state initialization
st.session_state.setdefault("vectors", None)
st.session_state.setdefault("final_documents", {})
st.session_state.setdefault("structured_data", {})  # {file_name: pandas.DataFrame or list for pptx}
st.session_state.setdefault("chat_folders", {"Default": []})  # {folder_name: [{chat_id, chat_name, messages: [{message, content, timestamp, retrieved_docs, pdf_data, is_table}]}]}
st.session_state.setdefault("uploaded_files", {})
st.session_state.setdefault("current_folder", "Default")
st.session_state.setdefault("current_chat_id", None)

# === [ Web Scraping for TSEC Website ] ===
def scrape_tsec_website():
    try:
        response = requests.get("https://tsec.edu/", timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        
        content = ""
        for element in soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6", "div"]):
            text = element.get_text(strip=True)
            if text:
                content += text + "\n"
        
        content = re.sub(r"\s+", " ", content).strip()
        if not content:
            logger.warning("No content scraped from TSEC website")
            return []
        
        doc = Document(page_content=content, metadata={"file_name": "tsec_website", "section": "website_info"})
        splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
        docs = splitter.split_documents([doc])
        for d in docs:
            if not isinstance(d, Document):
                logger.error(f"Non-Document object from website splitter: {type(d)}")
                return []
        return docs
    except Exception as e:
        logger.error(f"Failed to scrape TSEC website: {e}")
        return []

# === [ Process Training Data ] ===
def process_training_file(file_path, file_name):
    try:
        if file_path.endswith(".pdf"):
            docs = PyPDFLoader(file_path).load()
            return [d for d in docs if isinstance(d, Document)], None
        elif file_path.endswith(".docx"):
            doc = docx.Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            if not text.strip():
                logger.warning(f"No text extracted from DOCX: {file_name}")
                return [], None
            return [Document(page_content=text, metadata={"file_name": file_name, "file_type": "docx"})], None
        elif file_path.endswith((".jpg", ".png")):
            img = Image.open(file_path).convert("L")  # Convert to grayscale
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(2.0)  # Increase contrast
            text = pytesseract.image_to_string(img)
            if not text.strip():
                logger.warning(f"No text extracted from image: {file_name}")
                return [], None
            metadata = {
                "file_name": file_name,
                "file_type": "image",
                "resolution": f"{img.width}x{img.height}",
                "size_bytes": os.path.getsize(file_path)
            }
            return [Document(page_content=text, metadata=metadata)], None
        elif file_path.endswith(".wav"):
            recognizer = sr.Recognizer()
            with sr.AudioFile(file_path) as source:
                audio = recognizer.record(source)
                text = recognizer.recognize_google(audio)
            if not text.strip():
                logger.warning(f"No text extracted from WAV: {file_name}")
                return [], None
            return [Document(page_content=text, metadata={"file_name": file_name, "file_type": "audio"})], None
        elif file_path.endswith((".csv", ".xlsx", ".xls")):
            try:
                if file_path.endswith(".csv"):
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)
                text_content = df.to_json(orient="records", lines=True)
                doc = Document(
                    page_content=text_content,
                    metadata={"file_name": file_name, "file_type": "structured", "columns": df.columns.tolist()}
                )
                splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
                docs = splitter.split_documents([doc])
                return [d for d in docs if isinstance(d, Document)], df
            except Exception as e:
                logger.error(f"Failed to process Excel/CSV {file_name}: {e}")
                return [], None
        elif file_path.endswith((".ppt", ".pptx")):
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                slide_content = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content += shape.text + "\n"
                if slide_content.strip():
                    slides_text.append(slide_content.strip())
            if not slides_text:
                logger.warning(f"No text extracted from PowerPoint: {file_name}")
                return [], None
            text_content = "\n\n".join(slides_text)
            doc = Document(
                page_content=text_content,
                metadata={"file_name": file_name, "file_type": "powerpoint", "slide_count": len(slides_text)}
            )
            splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
            docs = splitter.split_documents([doc])
            return [d for d in docs if isinstance(d, Document)], slides_text
        logger.warning(f"Unsupported file type: {file_name}")
        return [], None
    except FileNotFoundError:
        logger.error(f"File not found: {file_path}")
        return [], None
    except Exception as e:
        logger.error(f"Error processing {file_path}: {e}")
        return [], None

def load_training_data():
    # Always rebuild FAISS index to ensure consistency
    if Path(VECTOR_STORE_PATH).exists():
        try:
            shutil.rmtree(VECTOR_STORE_PATH)
            logger.info("Cleared existing FAISS index for rebuilding")
        except Exception as e:
            logger.error(f"Failed to clear FAISS index: {e}")
    
    # Clear final_documents to start fresh
    st.session_state.final_documents = {}
    st.session_state.structured_data = {}
    st.session_state.vectors = None

    for file_name in os.listdir(TRAINING_DATA_PATH):
        file_path = os.path.join(TRAINING_DATA_PATH, file_name)
        docs, data = process_training_file(file_path, file_name)
        if not docs:
            continue
        final_docs = []
        for doc in docs:
            if not isinstance(doc, Document):
                logger.error(f"Non-Document object in {file_name}: {type(doc)}")
                continue
            content = doc.page_content.lower()
            doc.metadata["file_name"] = file_name
            internship_keywords = ["internship", "intern", "experience", "training"]
            project_keywords = ["project", "projects"]
            education_keywords = ["education"]
            skills_keywords = ["skills"]
            if any(k in content for k in internship_keywords):
                doc.metadata["section"] = "internship_or_experience"
            elif any(k in content for k in education_keywords):
                doc.metadata["section"] = "education"
            elif any(k in content for k in skills_keywords):
                doc.metadata["section"] = "skills"
            elif any(k in content for k in project_keywords):
                doc.metadata["section"] = "projects"
            elif doc.metadata.get("file_type") == "structured":
                doc.metadata["section"] = "structured_data"
            elif doc.metadata.get("file_type") == "powerpoint":
                doc.metadata["section"] = "powerpoint"
            elif doc.metadata.get("file_type") == "image":
                doc.metadata["section"] = "image"
            else:
                doc.metadata["section"] = "other"
            final_docs.append(doc)

        if final_docs:
            st.session_state.final_documents[file_name] = final_docs
            if data is not None:
                st.session_state.structured_data[file_name] = data
            if st.session_state.vectors is None:
                st.session_state.vectors = FAISS.from_documents(final_docs, embeddings)
            else:
                st.session_state.vectors.add_documents(final_docs)

    if "tsec_website" not in st.session_state.final_documents:
        website_docs = scrape_tsec_website()
        if website_docs:
            valid_docs = [d for d in website_docs if isinstance(d, Document)]
            if valid_docs:
                st.session_state.final_documents["tsec_website"] = valid_docs
                if st.session_state.vectors is None:
                    st.session_state.vectors = FAISS.from_documents(valid_docs, embeddings)
                else:
                    st.session_state.vectors.add_documents(valid_docs)
            else:
                logger.error("No valid Document objects from TSEC website")

    if st.session_state.vectors:
        st.session_state.vectors.save_local(VECTOR_STORE_PATH)
        logger.info("Training data and website content embedded and indexed successfully!")

# Initialize vector store by always rebuilding
load_training_data()

# === [ Generate Chat Name ] ===
def generate_chat_name(question):
    prompt_template = ChatPromptTemplate.from_template(
        "Generate a concise chat name (5-8 words) based on this question: {question}"
    )
    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = llm.invoke(prompt_template.format_prompt(question=question).to_messages())
            name = response.content.strip()[:50]
            return name or f"Chat {datetime.now().strftime('%H:%M')}"
        except groq.AuthenticationError as e:
            st.error(f"Authentication failed in chat name generation: {e}. Please check your Groq API key.")
            return f"Chat {datetime.now().strftime('%H:%M')}"
        except groq.RateLimitError as e:
            try:
                wait_time_str = e.response.json()['error']['message'].split('Please try again in ')[1].split(' ')[0]
                wait_time = float(wait_time_str.replace('m', '*60').replace('s', ''))
            except (IndexError, ValueError):
                wait_time = 5
            st.error(f"Rate limit exceeded for Groq API. Attempt {attempt + 1}/{max_retries}. Retrying in {wait_time} seconds.")
            if attempt < max_retries - 1:
                time.sleep(wait_time + 1)
            else:
                return f"Chat {datetime.now().strftime('%H:%M')}"
        except Exception as e:
            logger.error(f"Error generating chat name: {e}")
            return f"Chat {datetime.now().strftime('%H:%M')}"

# === [ Voice Command Processing ] ===
def process_voice_command(text):
    text_lower = text.lower()
    commands = {
        "new chat": lambda: create_new_chat(),
        "delete chat": lambda: delete_current_chat(),
        "summarize file": lambda: trigger_summary(),
        "list internships": lambda: return_query("List all internships"),
        "clear data": lambda: clear_all_data()
    }
    for cmd, action in commands.items():
        if cmd in text_lower:
            return action()
    return text  # Return text as query if no command matches

def create_new_chat():
    chat_id = str(uuid.uuid4())
    chat_entry = {
        'chat_id': chat_id,
        'chat_name': f"New Chat {datetime.now().strftime('%H:%M')}",
        'messages': [],
        'chat_name_set': False
    }
    st.session_state.chat_folders[st.session_state.current_folder].append(chat_entry)
    st.session_state.current_chat_id = chat_id
    st.success("✅ New chat started!")
    return None

def delete_current_chat():
    if st.session_state.current_chat_id:
        st.session_state.chat_folders[st.session_state.current_folder] = [
            chat for chat in st.session_state.chat_folders[st.session_state.current_folder]
            if chat['chat_id'] != st.session_state.current_chat_id
        ]
        st.session_state.current_chat_id = None
        st.success("✅ Current chat deleted!")
    return None

def trigger_summary():
    if st.session_state.uploaded_files:
        st.session_state["trigger_summary"] = True
    return None

def clear_all_data():
    st.session_state.vectors = None
    st.session_state.final_documents = {}
    st.session_state.structured_data = {}
    st.session_state.chat_folders = {"Default": []}
    st.session_state.uploaded_files = {}
    st.session_state.current_folder = "Default"
    st.session_state.current_chat_id = None
    for file_name in os.listdir(STORAGE_PATH):
        try:
            os.remove(os.path.join(STORAGE_PATH, file_name))
        except Exception as e:
            st.warning(f"Failed to delete {file_name}: {e}")
    if os.path.exists(VECTOR_STORE_PATH):
        try:
            shutil.rmtree(VECTOR_STORE_PATH)
            logger.info("Cleared FAISS index directory")
        except Exception as e:
            st.warning(f"Failed to remove directory {VECTOR_STORE_PATH}: {e}")
    st.success("✅ All data cleared!")
    load_training_data()
    return None

def return_query(query):
    return query

# === [ Sidebar ] ===
with st.sidebar:
    st.header("📜 Chats", divider="blue")
    
    with st.expander("🗂 Folders", expanded=True):
        folder_options = list(st.session_state.chat_folders.keys())
        selected_folder = st.selectbox("Select Folder", options=folder_options, key="folder_select")
        st.session_state.current_folder = selected_folder
        
        col1, col2 = st.columns(2)
        with col1:
            new_folder_name = st.text_input("New Folder", label_visibility="collapsed", key="new_folder_input")
        with col2:
            if st.button("➕", key="add_folder"):
                if new_folder_name and new_folder_name.strip() not in folder_options:
                    st.session_state.chat_folders[new_folder_name.strip()] = []
                    st.session_state.current_folder = new_folder_name.strip()
                    st.success(f"✅ Folder '{new_folder_name.strip()}' created!")
                else:
                    st.error("Enter a unique folder name.")
        
        if selected_folder != "Default":
            col1, col2 = st.columns(2)
            with col1:
                rename_folder = st.text_input("Rename Folder", value=selected_folder, key=f"rename_folder_{selected_folder}")
                if st.button("Save", key=f"save_folder_{selected_folder}"):
                    if rename_folder.strip() and rename_folder.strip() not in folder_options:
                        st.session_state.chat_folders[rename_folder.strip()] = st.session_state.chat_folders.pop(selected_folder)
                        st.session_state.current_folder = rename_folder.strip()
                        st.success(f"✅ Folder renamed to '{rename_folder.strip()}'!")
                    else:
                        st.error("Enter a unique folder name.")
            with col2:
                if st.button("🗑️", key=f"delete_folder_{selected_folder}"):
                    if not st.session_state.chat_folders[selected_folder]:
                        st.session_state.chat_folders.pop(selected_folder)
                        st.session_state.current_folder = "Default"
                        st.success(f"✅ Folder '{selected_folder}' deleted!")
                    else:
                        st.error("Cannot delete folder with chats.")
    
    for folder in folder_options:
        with st.expander(f"📁 {folder}", expanded=folder == st.session_state.current_folder):
            if not st.session_state.chat_folders[folder]:
                st.write("No chats in this folder.")
            for chat in st.session_state.chat_folders[folder]:
                chat_id = chat.get('chat_id', str(uuid.uuid4()))
                if 'chat_id' not in chat:
                    chat['chat_id'] = chat_id
                
                col1, col2 = st.columns([0.8, 0.2])
                with col1:
                    if st.button(chat.get('chat_name', 'Unnamed Chat'), key=f"select_{chat_id}_{folder}"):
                        st.session_state.current_folder = folder
                        st.session_state.current_chat_id = chat_id
                with col2:
                    if st.button("⋮", key=f"menu_{chat_id}_{folder}"):
                        st.session_state[f"show_menu_{chat_id}"] = not st.session_state.get(f"show_menu_{chat_id}", False)
                
                if st.session_state.get(f"show_menu_{chat_id}", False):
                    with st.container():
                        rename_input = st.text_input("Rename Chat", value=chat.get('chat_name', ''), key=f"rename_{chat_id}_{folder}")
                        if st.button("Save Name", key=f"save_rename_{chat_id}_{folder}"):
                            chat['chat_name'] = rename_input.strip()
                            st.session_state[f"show_menu_{chat_id}"] = False
                            st.success("✅ Chat renamed!")
                        
                        if st.button("Delete Chat", key=f"delete_{chat_id}_{folder}"):
                            st.session_state.chat_folders[folder] = [
                                c for c in st.session_state.chat_folders[folder] if c['chat_id'] != chat_id
                            ]
                            if st.session_state.current_chat_id == chat_id:
                                st.session_state.current_chat_id = None
                            st.session_state[f"show_menu_{chat_id}"] = False
                            st.success("✅ Chat deleted!")
                        
                        target_folder = st.selectbox("Move to Folder", options=[f for f in folder_options if f != folder], key=f"move_{chat_id}_{folder}")
                        if st.button("Move", key=f"move_btn_{chat_id}_{folder}"):
                            chat_copy = chat.copy()
                            st.session_state.chat_folders[target_folder].append(chat_copy)
                            st.session_state.chat_folders[folder] = [
                                c for c in st.session_state.chat_folders[folder] if c['chat_id'] != chat_id
                            ]
                            if st.session_state.current_chat_id == chat_id:
                                st.session_state.current_folder = target_folder
                                st.session_state.current_chat_id = chat_id
                            st.session_state[f"show_menu_{chat_id}"] = False
                            st.success(f"✅ Chat moved to {target_folder}!")

    if st.button("🗗 New Chat", key="new_chat"):
        create_new_chat()

# === [ UI: Title, Upload, & Tabs ] ===
st.title("🎓 College Assistant RAG Chatbot")

uploaded_files = st.file_uploader(
    "📤 Upload Files (PDF, DOCX, JPG, PNG, WAV, CSV, XLSX, XLS, PPT, PPTX)",
    type=["pdf", "docx", "jpg", "png", "wav", "csv", "xlsx", "xls", "ppt", "pptx"],
    accept_multiple_files=True
)

if st.button("🗑️ Clear All Data"):
    clear_all_data()

if uploaded_files:
    with st.spinner("🔄 Processing uploaded files..."):
        for uploaded_file in uploaded_files:
            file_path = os.path.join(STORAGE_PATH, uploaded_file.name)
            if uploaded_file.name not in st.session_state.uploaded_files:
                with open(file_path, "wb") as f:
                    f.write(uploaded_file.getvalue())
                st.session_state.uploaded_files[uploaded_file.name] = file_path

        for file_name, file_path in st.session_state.uploaded_files.items():
            if file_name not in st.session_state.final_documents:
                docs, data = process_training_file(file_path, file_name)
                if not docs:
                    continue
                final_docs = []
                for doc in docs:
                    if not isinstance(doc, Document):
                        logger.error(f"Non-Document object in uploaded {file_name}: {type(doc)}")
                        continue
                    content = doc.page_content.lower()
                    doc.metadata["file_name"] = file_name
                    internship_keywords = ["internship", "intern", "experience", "training"]
                    project_keywords = ["project", "projects"]
                    education_keywords = ["education"]
                    skills_keywords = ["skills"]
                    if any(k in content for k in internship_keywords):
                        doc.metadata["section"] = "internship_or_experience"
                    elif any(k in content for k in education_keywords):
                        doc.metadata["section"] = "education"
                    elif any(k in content for k in skills_keywords):
                        doc.metadata["section"] = "skills"
                    elif any(k in content for k in project_keywords):
                        doc.metadata["section"] = "projects"
                    elif doc.metadata.get("file_type") == "structured":
                        doc.metadata["section"] = "structured_data"
                    elif doc.metadata.get("file_type") == "powerpoint":
                        doc.metadata["section"] = "powerpoint"
                    elif doc.metadata.get("file_type") == "image":
                        doc.metadata["section"] = "image"
                    else:
                        doc.metadata["section"] = "other"
                    final_docs.append(doc)

                if final_docs:
                    st.session_state.final_documents[file_name] = final_docs
                    if data is not None:
                        st.session_state.structured_data[file_name] = data
                    if st.session_state.vectors is None:
                        st.session_state.vectors = FAISS.from_documents(final_docs, embeddings)
                    else:
                        st.session_state.vectors.add_documents(final_docs)

        if st.session_state.vectors:
            st.session_state.vectors.save_local(VECTOR_STORE_PATH)
            st.success("✅ Files embedded and indexed successfully!")

st.write("📂 **Uploaded Files**:")
for file_name in st.session_state.uploaded_files:
    st.write(f"📔 {file_name}")

# === [ Summary Generation ] ===
def generate_summary(file_name):
    docs = st.session_state.final_documents.get(file_name, [])
    if not docs:
        return "No content available to summarize."
    valid_docs = [d for d in docs if isinstance(d, Document)]
    if not valid_docs:
        logger.error(f"No valid Document objects for summary in {file_name}")
        return "Invalid document format."
    if file_name in st.session_state.structured_data:
        data = st.session_state.structured_data[file_name]
        if isinstance(data, pd.DataFrame):
            context = data.to_json(orient="records", lines=True)
        else:  # PowerPoint
            context = "\n".join(data)
    else:
        context = "\n".join([doc.page_content for doc in valid_docs])
    prompt_template = ChatPromptTemplate.from_template(
        """Summarize the following document content concisely in up to 200 words, focusing on key information such as main topics, purpose, or significant details. For structured data, highlight key columns and insights. For PowerPoint, summarize slide content. Provide a clear and structured summary.

        Document Content:
        {context}

        Summary:"""
    )
    document_chain = create_stuff_documents_chain(llm, prompt_template)
    max_retries = 3
    for attempt in range(max_retries):
        try:
            summary = document_chain.invoke({'context': context, 'input': ''})
            return summary
        except groq.AuthenticationError as e:
            st.error(f"Authentication failed in summary generation: {e}. Please check your Groq API key.")
            return None
        except groq.RateLimitError as e:
            try:
                wait_time_str = e.response.json()['error']['message'].split('Please try again in ')[1].split(' ')[0]
                wait_time = float(wait_time_str.replace('m', '*60').replace('s', ''))
            except (IndexError, ValueError):
                wait_time = 5
            st.error(f"Rate limit exceeded for Groq API. Attempt {attempt + 1}/{max_retries}. Retrying in {wait_time} seconds.")
            if attempt < max_retries - 1:
                time.sleep(wait_time + 1)
            else:
                return None
        except Exception as e:
            st.error(f"Error generating summary: {e}")
            return None

if st.session_state.uploaded_files:
    st.write("📝 **Generate Summary**")
    selected_file = st.selectbox("Select File to Summarize:", [""] + list(st.session_state.uploaded_files.keys()), key="summary_select")
    if (selected_file and st.button("Generate Summary", key="generate_summary")) or st.session_state.get("trigger_summary"):
        with st.spinner("📄 Generating summary..."):
            summary = generate_summary(selected_file)
            if summary:
                st.markdown(f"### Summary for {selected_file}")
                st.write(summary)
                html_content = f"""
                <h1>Summary for {selected_file}</h1>
                <p>{summary}</p>
                """
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                    temp_file_path = temp_file.name
                try:
                    pdfkit.from_string(html_content, temp_file_path, configuration=pdfkit_config)
                    with open(temp_file_path, "rb") as f:
                        pdf_data = f.read()
                finally:
                    try:
                        os.unlink(temp_file_path)
                    except PermissionError:
                        st.warning(f"Could not delete temporary file {temp_file_path}. Please close any programs using it and delete it manually.")
                st.download_button(
                    label="📥 Download Summary as PDF",
                    data=pdf_data,
                    file_name=f"summary_{selected_file}.pdf",
                    mime="application/pdf",
                    key="download_summary"
                )
        st.session_state["trigger_summary"] = False

# === [ Search Logic ] ===
def extract_name(query):
    possible_names = list(st.session_state.uploaded_files.keys())
    name_match = re.findall(r"\b[A-Z][a-z]+(?:\s[A-Z][a-z]+)?\b", query)
    for token in name_match:
        for file_name in possible_names:
            if token.lower() in file_name.lower():
                return token, file_name
    return None, None

def hybrid_search(query):
    matched_docs = []
    query_lower = query.lower()
    name, associated_file = extract_name(query)
    file_scope = [associated_file] if associated_file else list(st.session_state.final_documents.keys())

    if st.session_state.vectors:
        retriever = st.session_state.vectors.as_retriever(search_kwargs={"k": 10})
        retrieved_docs = retriever.invoke(query)
        logger.info(f"Retriever output for query '{query}': {[type(doc) for doc in retrieved_docs]}")
        
        for doc in retrieved_docs:
            if isinstance(doc, Document):
                matched_docs.append(doc)
            else:
                logger.warning(f"Skipping non-Document object: {type(doc)} - {doc}")

    return matched_docs

# === [ Table Detection and Generation ] ===
def needs_table(query):
    table_keywords = ["list", "table", "compare", "summary", "data", "rows", "columns", "slide", "presentation"]
    return any(keyword in query.lower() for keyword in table_keywords)

def generate_table_content(query, docs):
    if not docs or not any(doc.metadata.get("file_type") in ["structured", "powerpoint"] for doc in docs):
        return None, None
    valid_docs = [d for d in docs if isinstance(d, Document)]
    if not valid_docs:
        logger.error("No valid Document objects for table generation")
        return None, None
    relevant_files = set(doc.metadata.get("file_name") for doc in valid_docs if doc.metadata.get("file_type") in ["structured", "powerpoint"])
    table_data = []
    query_lower = query.lower()
    for file_name in relevant_files:
        data = st.session_state.structured_data.get(file_name)
        if data is None:
            continue
        if isinstance(data, pd.DataFrame):
            if "list" in query_lower or "table" in query_lower:
                table_data.extend(data.head(5).to_dict(orient="records"))
            elif "compare" in query_lower:
                columns = [col for col in data.columns if any(col.lower() in query_lower for col in data.columns)]
                if columns:
                    table_data.extend(data[columns].head(5).to_dict(orient="records"))
            elif "summary" in query_lower:
                summary = data.describe().to_dict()
                table_data = [{"Statistic": stat, **{col: summary[col][stat] for col in summary}} for stat in summary[list(summary.keys())[0]]]
        elif isinstance(data, list):  # PowerPoint
            table_data.extend([{"Slide": i + 1, "Content": content[:100] + "..."} for i, content in enumerate(data)])
    if table_data:
        columns = list(table_data[0].keys())
        return table_data, columns
    return None, None

# === [ PDF Generation ] ===
def generate_answer_pdf(message, content, timestamp, table_data=None, table_columns=None):
    html_content = f"""
    <html>
    <head><style>body {{ font-family: Arial; color: white; background-color: #1e1e1e; }}</style></head>
    <body>
    <h1>Answer to Question</h1>
    <h3>Question: {message}</h3>
    <p><strong>Answer:</strong><br>{content}</p>
    """
    if table_data and table_columns:
        html_content += "<h3>Table Data</h3><table border='1' style='border-collapse: collapse;'>"
        html_content += "<tr>" + "".join(f"<th>{col}</th>" for col in table_columns) + "</tr>"
        for row in table_data:
            html_content += "<tr>" + "".join(f"<td>{row.get(col, '')}</td>" for col in table_columns) + "</tr>"
        html_content += "</table>"
    html_content += f"<p style='text-align: right; color: #FFFFFF;'>{timestamp}</p>"
    html_content += "</body></html>"
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
        temp_file_path = temp_file.name
    try:
        pdfkit.from_string(html_content, temp_file_path, configuration=pdfkit_config)
        with open(temp_file_path, "rb") as f:
            pdf_data = f.read()
    finally:
        try:
            os.unlink(temp_file_path)
        except PermissionError:
            st.warning(f"Could not delete temporary file {temp_file_path}. Please close any programs using it and delete it manually.")
    return pdf_data

def generate_conversation_pdf(chat):
    html_content = f"""
    <html>
    <head><style>body {{ font-family: Arial; color: white; background-color: #1e1e1e; }}</style></head>
    <body>
    <h1>College Assistant RAG Chatbot Conversation</h1>
    <h2>{chat.get('chat_name', 'Unnamed Chat')}</h2>
    <hr>
    """
    for turn in chat['messages']:
        html_content += f"""
        <h3>You: {turn['message']}</h3>
        <p><strong>Answer:</strong><br>{turn['content']}</p>
        """
        if turn.get('is_table') and turn.get('table_data') and turn.get('table_columns'):
            html_content += "<h4>Table Data</h4><table border='1' style='border-collapse: collapse;'>"
            html_content += "<tr>" + "".join(f"<th>{col}</th>" for col in turn['table_columns']) + "</tr>"
            for row in turn['table_data']:
                html_content += "<tr>" + "".join(f"<td>{row.get(col, '')}</td>" for col in turn['table_columns']) + "</tr>"
            html_content += "</table>"
        html_content += f"<p style='text-align: right; color: #FFFFFF;'>{turn['timestamp']}</p><hr>"
    html_content += "</body></html>"
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
        temp_file_path = temp_file.name
    try:
        pdfkit.from_string(html_content, temp_file_path, configuration=pdfkit_config)
        with open(temp_file_path, "rb") as f:
            pdf_data = f.read()
    finally:
        try:
            os.unlink(temp_file_path)
        except PermissionError:
            st.warning(f"Could not delete temporary file {temp_file_path}. Please close any programs using it and delete it manually.")
    return pdf_data

# === [ Chat Tabs ] ===
if st.session_state.chat_folders.get(st.session_state.current_folder):
    tabs = [chat['chat_name'] for chat in st.session_state.chat_folders[st.session_state.current_folder]]
    tab_cols = st.columns(min(len(tabs) + 1, 6))
    for idx, chat in enumerate(st.session_state.chat_folders[st.session_state.current_folder]):
        if idx < 5:
            with tab_cols[idx]:
                is_active = chat['chat_id'] == st.session_state.current_chat_id
                style = "background-color: #007bff; color: white;" if is_active else ""
                if st.button(chat['chat_name'], key=f"tab_{chat['chat_id']}", help=chat['chat_name']):
                    st.session_state.current_chat_id = chat['chat_id']
    with tab_cols[-1]:
        if st.button("➕ New Chat", key="new_chat_tab"):
            create_new_chat()

# === [ Input Bar with Mic ] ===
col1, col2 = st.columns([0.9, 0.1])
with col1:
    prompt = st.text_input("❓ Ask a question:", value=st.session_state.get("text_input", ""), key="question_input", placeholder="Type or speak your question")
with col2:
    if st.button("🎙️", key="voice_input"):
        recognizer = sr.Recognizer()
        timeout = 7  # Increased timeout
        try:
            with sr.Microphone() as source:
                st.info("Listening... Speak clearly...")
                progress = st.progress(0)
                recognizer.adjust_for_ambient_noise(source, duration=2.0)  # Enhanced noise cancellation
                start_time = time.time()
                while time.time() - start_time < timeout:
                    progress.progress(min(1.0, (time.time() - start_time) / timeout))
                    try:
                        audio = recognizer.listen(source, timeout=1.0, phrase_time_limit=timeout)
                        break
                    except sr.WaitTimeoutError:
                        continue
                progress.empty()
                st.info("Processing audio...")
            try:
                text = recognizer.recognize_google(audio, language="en-US")
                processed_text = process_voice_command(text)
                if processed_text and isinstance(processed_text, str):
                    st.session_state["text_input"] = processed_text
                    st.success(f"Voice command processed: {processed_text}")
                    prompt = processed_text
                elif processed_text is None:
                    st.session_state["text_input"] = ""
            except sr.UnknownValueError:
                st.error("Could not understand audio. Please try again.")
                st.button("Retry Voice Input", key="retry_voice")
            except sr.RequestError as e:
                st.error(f"Speech recognition failed: {e}. Trying offline recognition...")
                try:
                    from vosk import Model, KaldiRecognizer
                    import wave
                    model = Model("vosk-model-small-en-us")  # Requires model files
                    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_wav:
                        with wave.open(temp_wav.name, "wb") as wf:
                            wf.setnchannels(1)
                            wf.setsampwidth(2)
                            wf.setframerate(16000)
                            wf.writeframes(audio.get_wav_data())
                        recognizer = KaldiRecognizer(model, 16000)
                        with open(temp_wav.name, "rb") as f:
                            recognizer.AcceptWaveform(f.read())
                        result = json.loads(recognizer.Result())
                        text = result.get("text", "")
                        if text:
                            processed_text = process_voice_command(text)
                            if processed_text and isinstance(processed_text, str):
                                st.session_state["text_input"] = processed_text
                                st.success(f"Offline voice processed: {processed_text}")
                                prompt = processed_text
                            elif processed_text is None:
                                st.session_state["text_input"] = ""
                        else:
                            st.error("Offline recognition failed.")
                            st.button("Retry Voice Input", key="retry_voice_offline")
                    os.unlink(temp_wav.name)
                except ImportError:
                    st.error("Offline speech recognition (vosk) not installed. Install with 'pip install vosk' and download a model.")
                except Exception as e:
                    st.error(f"Offline recognition error: {e}")
                    st.button("Retry Voice Input", key="retry_voice_error")
        except Exception as e:
            st.error(f"Microphone error: {e}. Ensure microphone is connected and permissions granted.")
            st.button("Retry Voice Input", key="retry_voice_mic")

# === [ Chat Interface ] ===
if prompt:
    if not st.session_state.current_chat_id:
        create_new_chat()
    
    with st.spinner("🔄 Generating response..."):
        results = hybrid_search(prompt)
        # Final safeguard: filter and log any remaining non-Document objects
        valid_results = []
        for doc in results:
            if not isinstance(doc, Document):
                logger.error(f"Non-Document object in final search results: {type(doc)} - {doc}")
                continue
            valid_results.append(doc)
        if len(results) != len(valid_results):
            logger.warning(f"Filtered {len(results) - len(valid_results)} non-Document objects from final search results")

        table_data, table_columns = None, None
        is_table = needs_table(prompt)
        if is_table:
            table_data, table_columns = generate_table_content(prompt, valid_results)

        if valid_results or is_table:
            prompt_template = ChatPromptTemplate.from_template(
                """**You are a helpful college assistant. ONLY use the context below to answer concisely and informatively. For questions requiring lists or comparisons, format as tables if possible. For PowerPoint or images, include slide or image content details.**
                Context:
                {context}
                Question:
                {input}
                Answer:"""
            )
            document_chain = create_stuff_documents_chain(llm, prompt_template)
            max_retries = 3
            response = None
            for attempt in range(max_retries):
                try:
                    context = "\n".join([doc.page_content for doc in valid_results])
                    if table_data:
                        context += f"\nTable Data (JSON):\n{json.dumps(table_data, default=str)}"
                    response = document_chain.invoke({'context': context, 'input': prompt})
                    break
                except groq.AuthenticationError as e:
                    st.error(f"Authentication failed: {e}. Please check your Groq API key in the .env file or Groq Console.")
                    response = None
                    break
                except groq.RateLimitError as e:
                    try:
                        wait_time_str = e.response.json()['error']['message'].split('Please try again in ')[1].split(' ')[0]
                        wait_time = float(wait_time_str.replace('m', '*60').replace('s', ''))
                    except (IndexError, ValueError):
                        wait_time = 5
                    st.error(f"Rate limit exceeded for Groq API. Attempt {attempt + 1}/{max_retries}. Retrying in {wait_time:.2f} seconds.")
                    if attempt < max_retries - 1:
                        time.sleep(wait_time + 1)
                    else:
                        response = None
                        break
                except Exception as e:
                    st.error(f"Error generating response: {e}")
                    response = None
                    break

            if response:
                timestamp = datetime.now().strftime("%H:%M:%S")
                pdf_data = generate_answer_pdf(prompt, response, timestamp, table_data, table_columns)
                message_entry = {
                    'message': prompt,
                    'content': response,
                    'timestamp': timestamp,
                    'retrieved_docs': valid_results[:10],
                    'pdf_data': pdf_data,
                    'is_table': is_table,
                    'table_data': table_data,
                    'table_columns': table_columns
                }
                
                for chat in st.session_state.chat_folders[st.session_state.current_folder]:
                    if chat['chat_id'] == st.session_state.current_chat_id:
                        if not chat['messages'] or not chat.get('chat_name_set', False):
                            chat['chat_name'] = generate_chat_name(prompt)
                            chat['chat_name_set'] = True
                        chat['messages'].append(message_entry)
                        break
                
                st.session_state["text_input"] = ""
            else:
                st.error("Failed to generate response. Please try again.")
                st.session_state["text_input"] = ""
        else:
            st.warning("No relevant documents found.")
            st.session_state["text_input"] = ""

# === [ Display Chat History ] ===
st.markdown("---")
if st.session_state.current_chat_id:
    for chat in st.session_state.chat_folders[st.session_state.current_folder]:
        if chat['chat_id'] == st.session_state.current_chat_id:
            st.markdown(f"### 💬 {chat['chat_name']}", unsafe_allow_html=True)
            if not chat['messages']:
                st.write("Start your conversation by asking a question!")
            else:
                for idx, turn in enumerate(chat['messages']):
                    st.markdown(f"<div style='margin-bottom: 10px;'><b>You:</b> {turn['message']}</div>", unsafe_allow_html=True)
                    st.markdown(
                        f"<div style='background-color: #1e1e1e; color: #FFFFFF; padding: 12px; border-radius: 8px; font-size: 16px; border: 1px solid #333333; margin-bottom: 5px;'><b>Answer:</b><br>{turn['content']}</div>",
                        unsafe_allow_html=True
                    )
                    if turn.get('is_table') and turn.get('table_data') and turn.get('table_columns'):
                        st.markdown("**Table Data**")
                        st.table(pd.DataFrame(turn['table_data'], columns=turn['table_columns']))
                    st.markdown(
                        f"<div style='text-align: right; font-size: 12px; color: #FFFFFF;'>{turn['timestamp']}</div>",
                        unsafe_allow_html=True
                    )
                    with st.expander(f"🔍 Debug: Retrieved Documents for '{turn['message']}'"):
                        st.write(f"📋 Retrieved {len(turn['retrieved_docs'])} documents")
                        for doc in turn['retrieved_docs']:
                            if isinstance(doc, Document):
                                st.write(f"**File**: {doc.metadata.get('file_name', 'Unknown')} | **Section**: {doc.metadata.get('section', 'None')}")
                                st.write(doc.page_content[:200] + "...")
                            else:
                                st.write(f"**Invalid Document**: Type {type(doc)} - {doc}")
                    pdf_data = turn.get('pdf_data') or generate_answer_pdf(
                        turn['message'], turn['content'], turn['timestamp'],
                        turn.get('table_data'), turn.get('table_columns')
                    )
                    st.download_button(
                        label="📥 Download Answer as PDF",
                        data=pdf_data,
                        file_name=f"answer_{turn['timestamp'].replace(':', '-')}.pdf",
                        mime="application/pdf",
                        key=f"download_answer_{chat['chat_id']}_{idx}"
                    )
                
                st.download_button(
                    label="📥 Export Conversation as PDF",
                    data=generate_conversation_pdf(chat),
                    file_name=f"conversation_{chat['chat_name'].replace(' ', '_')}.pdf",
                    mime="application/pdf",
                    key=f"download_conversation_{chat['chat_id']}"
                )
            break
else:
    st.write("Select or start a new chat to begin!")